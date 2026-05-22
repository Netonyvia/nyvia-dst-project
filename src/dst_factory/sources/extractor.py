from io import BytesIO
from pypdf import PdfReader
from pptx import Presentation

from dst_factory.drive.client import GoogleDriveClient
from dst_factory.sources.models import SourceDocument, SourceFile


GOOGLE_DOC_MIME_TYPE = "application/vnd.google-apps.document"
PDF_MIME_TYPE = "application/pdf"
PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

class DriveSourceExtractor:

    def __init__(self, drive_client: GoogleDriveClient | None = None) -> None:
        self.drive_client = drive_client or GoogleDriveClient()

    def list_folder_sources(self, folder_id: str) -> list[SourceFile]:
        files = self.drive_client.list_files_in_folder(folder_id)

        return [SourceFile(id=file["id"], name=file["name"], mime_type=file["mimeType"]) for file in files]


    def extract_source_document(self, source_file: SourceFile) -> SourceDocument:
        if source_file.mime_type == GOOGLE_DOC_MIME_TYPE:
            text = self.drive_client.export_google_doc_As_text(source_file.id)

        elif source_file.mime_type == PDF_MIME_TYPE:
            file_bytes = self.drive_client.service.files().get_media(fileId=source_file.id).execute()
            text = self.extract_pdf_text(file_bytes)

        elif source_file.mime_type == PPTX_MIME_TYPE:
            file_bytes = self.drive_client.service.files().get_media(fileId=source_file.id).execute()
            text = self.extract_pptx_text(file_bytes)
            
        else:
            text = "[Unsupported file type for text estraction]"

        return SourceDocument(source_file=source_file, text=text)
    
    def extract_folder_document(self, folder_id: str) -> list[SourceDocument]:
        source_files = self.list_folder_sources(folder_id)

        return [self.extract_source_document(source_file) for source_file in source_files]
    
    def extract_pdf_text(self, file_bytes: bytes) -> str:
        reader = PdfReader(BytesIO(file_bytes))
        
        pages_text = []
        for index, page in enumerate(reader.pages, start = 1):
            try:
                page_text = page.extract_text() or ""
                pages_text.append(f"\n--- PDF PAGE {index} ---\n{page_text}")
            except Exception as e:
                print(f"Error extracting text from page {index}: {e}")

        return "\n".join(pages_text)
    
    def extract_pptx_text(self, file_bytes: bytes) -> str:
        presentation = Presentation(BytesIO(file_bytes))
        
        slides_text = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            slides_text.append(f"\n--- PPTX SLIDE {index} ---\n" + "\n".join(slide_text))

        return "\n".join(slides_text)